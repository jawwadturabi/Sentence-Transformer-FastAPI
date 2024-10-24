# extractors.py

import os
import tempfile
import PyPDF2
from docx import Document
import openpyxl
from pptx import Presentation
import requests
from pydub import AudioSegment
import concurrent.futures

from utils import upload_pdf_images_to_s3, call_openai_to_extract_from_images

def extract_text_from_pptx(file_stream):
    prs = Presentation(file_stream)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_image(file_stream, database_document_id):
    """Extract text from image files and send them to OpenAI for processing"""
    try:
        # Uploading the image to S3 and getting back the URL
        print("Uploading image to S3...")
        image_urls = upload_images_to_s3(file_stream, database_document_id)
        print(f"Uploaded image to S3: {image_urls}")

        # Calling OpenAI to extract text from the image URL
        if image_urls:
            print("Sending image to OpenAI for text extraction...")
            full_text = call_openai_to_extract_from_images(image_urls)
            print(f"Extracted text from images: {full_text}")
        else:
            full_text = "No images were uploaded to S3."
        
        return full_text

    except Exception as e:
        print(f"Error extracting text from image: {e}")
        return ""


def extract_text_from_docx(file_stream):
    document = Document(file_stream)
    return "\n".join([paragraph.text for paragraph in document.paragraphs])

def extract_text_from_xlsx(file_stream):
    workbook = openpyxl.load_workbook(file_stream, data_only=True)
    text = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            text.append(" ".join([str(cell) if cell is not None else "" for cell in row]))
    return "\n".join(text)

def extract_text_from_pdf(file_stream, database_document_id):
    """Extract text from a PDF file and create images for each page if necessary"""
    try:
        # Try to extract the text using PyPDF2 first
        reader = PyPDF2.PdfReader(file_stream)
        full_text = ""
        # for page in reader.pages:
        #     page_text = page.extract_text()
        #     if page_text:  # If text is found
        #         full_text += page_text
        
        # If no text is found, convert pages to images and process with OpenAI
        if not full_text:
            print("No text found, uploading pages as images and sending to OpenAI.")
            image_urls = upload_pdf_images_to_s3(file_stream, database_document_id)
            print(f"Uploaded {len(image_urls)} images to S3.")
            full_text = call_openai_to_extract_from_images(image_urls)
            print(f"Extracted text from images: {full_text}")
        return full_text

    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return ""
        
def split_audio(file_path, segment_duration=30000):
    """Split audio into smaller segments of given duration (in milliseconds)."""
    audio = AudioSegment.from_wav(file_path)
    segments = []
    for i in range(0, len(audio), segment_duration):
        segment = audio[i:i+segment_duration]
        segments.append(segment)
    return segments

def extract_audio_from_video(file_stream):
    """Extract the audio from an MP4 video file and convert it to WAV format."""
    # Convert the file stream to a temporary MP4 file
    with tempfile.NamedTemporaryFile(suffix=".mp4", delete=False) as temp_video_file:
        temp_video_file.write(file_stream.read())
        temp_video_path = temp_video_file.name

    # Use pydub to extract the audio from the MP4 file
    audio = AudioSegment.from_file(temp_video_path, format="mp4")
    
    # Create a temporary WAV file to store the extracted audio
    with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as temp_audio_file:
        audio.export(temp_audio_file.name, format="wav")
        temp_audio_path = temp_audio_file.name

    # Remove the temporary MP4 file after extracting the audio
    os.remove(temp_video_path)
    
    return temp_audio_path

def extract_text_from_audio(file_stream, file_type):
    """Transcribe audio using Whisper via Hugging Face API with parallel processing."""
    
    if file_type == 'mp4':
        print("Extracting audio from MP4 video...")
        temp_audio_path = extract_audio_from_video(file_stream)
    else:
        print("Creating temporary audio file...")
        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as temp_audio_file:
            temp_audio_file.write(file_stream.read())
            temp_audio_path = temp_audio_file.name

    # Split the audio into X-second segments
    segments = split_audio(temp_audio_path, segment_duration=30000)  # 30 seconds = 30,000 ms
    
    # API config
    API_URL = "https://api-inference.huggingface.co/models/openai/whisper-large-v3-turbo"
    API_TOKEN = os.environ['HUGGINGFACE_API_TOKEN']  # API token stored in the environment variables
    headers = {"Authorization": f"Bearer {API_TOKEN}"}

    # Use ThreadPoolExecutor to parallelize the requests to Hugging Face API
    transcriptions = {}

    with concurrent.futures.ThreadPoolExecutor() as executor:
        # Submit transcription tasks to the executor
        future_to_segment = {executor.submit(transcribe_segment, segment, i, headers, API_URL): i for i, segment in enumerate(segments)}
        
        # Collect results as they complete
        for future in concurrent.futures.as_completed(future_to_segment):
            i = future_to_segment[future]
            try:
                segment_text = future.result()
                transcriptions[i] = segment_text
            except Exception as e:
                print(f"Error during transcription of segment {i+1}: {e}")

    print("Transcription done for all segments.")
    
    # Sort the transcriptions by the segment index and concatenate them
    full_transcription = " ".join([transcriptions[i] for i in sorted(transcriptions)])

    # Delete the temporary audio file after use
    os.remove(temp_audio_path)
    
    return full_transcription

def transcribe_segment(segment, i, headers, API_URL):
    """Helper function to transcribe a single audio segment using Hugging Face API."""
    # Save the segment to a temporary file
    with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as temp_segment_file:
        segment.export(temp_segment_file.name, format="wav")
        segment_path = temp_segment_file.name

    # Read the audio segment
    with open(segment_path, "rb") as audio_file:
        audio_data = audio_file.read()

    # Send the request to the Hugging Face Whisper API
    print(f"Sending segment {i+1} to Hugging Face API...")
    response = requests.post(API_URL, headers=headers, data=audio_data)

    # Delete the temporary segment file after use
    os.remove(segment_path)

    # Check if the request was successful
    if response.status_code == 200:
        result = response.json()
        return result['text']
    else:
        raise Exception(f"Failed to transcribe segment {i+1}: {response.status_code}, {response.text}")